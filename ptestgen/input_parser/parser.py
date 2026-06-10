import os
from typing import List, Dict, Any, Tuple
import logging

# Placeholder for future PDF/DOCX/PPTX parsing libraries
# import pypdf
# import docx
# from pptx import Presentation
# from PIL import Image

# Import parsing libraries
try:
    import pypdf
except ImportError:
    pypdf = None
    logging.warning("`pypdf` library not found. Install with `pip install pypdf` to enable PDF parsing.")

try:
    import docx
except ImportError:
    docx = None
    logging.warning("`python-docx` library not found. Install with `pip install python-docx` to enable DOCX parsing.")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    logging.warning("`python-pptx` library not found. Install with `pip install python-pptx` to enable PPTX parsing.")

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None
    logging.warning("`striprtf` library not found. Install with `pip install striprtf` to enable RTF parsing.")

# PIL (Pillow) is still needed for image handling later, but not directly used here yet
try:
    from PIL import Image
except ImportError:
    Image = None
    # No warning needed here yet, as it's only for explicit image inputs for now

try:
    from nbconvert import MarkdownExporter
except ImportError:
    MarkdownExporter = None
    logging.warning("`nbconvert` library not found. Install with `pip install nbconvert` to enable IPYNB parsing.")


def parse_input_material(file_path: str) -> Tuple[str, list]:
    """
    Parses the input material file and returns its text content.
    Optionally extracts images found within the document (functionality not implemented yet).

    Args:
        file_path (str): Path to the input document.

    Returns:
        Tuple[str, list]: A tuple containing:
            - The extracted text content (str).
            - A list of extracted image references (currently placeholders or file paths) (list).
            Returns ("", []) if parsing fails or format is unsupported.

    Raises:
        FileNotFoundError: If the input file does not exist.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")

    _, extension = os.path.splitext(file_path)
    extension = extension.lower()
    text_content = ""
    image_references = [] # To hold references to extracted images later

    logging.info(f"Attempting to parse '{file_path}' (extension: {extension})")

    try:
        if extension == ".txt" or extension == ".md": # Treat Markdown as plain text for now
            with open(file_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
            logging.info(f"Successfully parsed {extension} file as plain text.")

        elif extension == ".pdf":
            if pypdf:
                try:
                    reader = pypdf.PdfReader(file_path)
                    text_parts = []
                    for i, page in enumerate(reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                             text_parts.append(page_text)
                    text_content = "\n".join(text_parts)
                    logging.info(f"Successfully parsed PDF file using pypdf. Extracted {len(text_content)} characters.")
                except Exception as e:
                    logging.error(f"Error parsing PDF file '{file_path}' with pypdf: {e}", exc_info=True)
            else:
                logging.warning(f"Skipping PDF parsing for '{file_path}' as pypdf library is not available.")

        elif extension == ".docx":
            if docx:
                try:
                    document = docx.Document(file_path)
                    text_parts = [para.text for para in document.paragraphs if para.text]
                    # Add text from tables (basic implementation)
                    for table in document.tables:
                         for row in table.rows:
                             for cell in row.cells:
                                 cell_text = cell.text.strip()
                                 if cell_text:
                                     text_parts.append(cell_text) # Append cell text as separate paragraph
                    text_content = "\n".join(text_parts)
                    logging.info(f"Successfully parsed DOCX file using python-docx. Extracted {len(text_content)} characters.")
                except Exception as e:
                     logging.error(f"Error parsing DOCX file '{file_path}' with python-docx: {e}", exc_info=True)
            else:
                logging.warning(f"Skipping DOCX parsing for '{file_path}' as python-docx library is not available.")

        elif extension == ".pptx":
            if Presentation:
                try:
                    prs = Presentation(file_path)
                    text_parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                shape_text = shape.text.strip()
                                if shape_text:
                                     text_parts.append(shape_text)
                            # --- Placeholder for Table Text ---
                            if shape.has_table:
                                table = shape.table
                                for row in table.rows:
                                    for cell in row.cells:
                                         cell_text = cell.text_frame.text.strip()
                                         if cell_text:
                                             text_parts.append(cell_text)
                    text_content = "\n".join(text_parts)
                    logging.info(f"Successfully parsed PPTX file using python-pptx. Extracted {len(text_content)} characters.")
                except Exception as e:
                    logging.error(f"Error parsing PPTX file '{file_path}' with python-pptx: {e}", exc_info=True)
            else:
                logging.warning(f"Skipping PPTX parsing for '{file_path}' as python-pptx library is not available.")

        elif extension == ".rtf":
            if rtf_to_text:
                 try:
                     with open(file_path, 'r', encoding='ascii', errors='ignore') as f: # RTF often uses extended ASCII
                         rtf_content = f.read()
                     text_content = rtf_to_text(rtf_content)
                     logging.info(f"Successfully parsed RTF file using striprtf. Extracted {len(text_content)} characters.")
                 except Exception as e:
                     logging.error(f"Error parsing RTF file '{file_path}' with striprtf: {e}", exc_info=True)
            else:
                 logging.warning(f"Skipping RTF parsing for '{file_path}' as striprtf library is not available.")

        elif extension == ".ipynb":
            if MarkdownExporter:
                try:
                    md_exporter = MarkdownExporter()
                    (text_content, resources) = md_exporter.from_filename(file_path)
                    logging.info(f"Successfully parsed IPYNB file by converting to Markdown. Extracted {len(text_content)} characters.")
                except Exception as e:
                    logging.error(f"Error parsing IPYNB file '{file_path}' with nbconvert: {e}", exc_info=True)
            else:
                logging.warning(f"Skipping IPYNB parsing for '{file_path}' as nbconvert library is not available.")
        else:
            logging.warning(f"Unsupported file format: {extension} for file '{file_path}'. Cannot extract text.")
            # Consider trying to read as plain text as a fallback?
            # try:
            #     with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            #         text_content = f.read()
            #     logging.info(f"Read unsupported file {extension} as plain text (might be garbled).")
            # except Exception as e:
            #      logging.error(f"Could not read unsupported file '{file_path}' as plain text: {e}")

    except Exception as e:
         logging.error(f"An unexpected error occurred during parsing of '{file_path}': {e}", exc_info=True)
         return "", [] # Return empty on major failure

    return text_content.strip(), image_references

def parse_image_input(image_path: str) -> Any:
    """
    Parses an input image file provided directly.
    This is distinct from extracting images embedded within documents.

    Args:
        image_path (str): Path to the image file.

    Returns:
        Any: Placeholder representing the loaded image data (currently the path).
             Returns None if parsing fails or format is unsupported.

    Raises:
        FileNotFoundError: If the image file does not exist.
        ValueError: If the image format is not supported.
    """
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image file not found: {image_path}")

    _, extension = os.path.splitext(image_path)
    supported_extensions = ['.png', '.jpg', '.jpeg', '.bmp', '.gif'] # Common vision model supported types
    if extension.lower() not in supported_extensions:
        raise ValueError(f"Unsupported image format for direct input: {extension}")

    # Check if Pillow is available for potential future validation/processing
    if Image:
        try:
            # Try opening the image to validate it
            with Image.open(image_path) as img:
                 img.verify() # Verify image header without loading full data
            logging.debug(f"Validated image file: {image_path}")
            # For now, just return the path as the 'parsed' data
            return image_path
        except Exception as e:
            logging.error(f"Failed to validate image file '{image_path}' using Pillow: {e}", exc_info=True)
            # Decide if we should raise ValueError or return None/path anyway
            # Let's raise for now, as it indicates a problem with the file
            raise ValueError(f"Image file '{image_path}' seems corrupted or invalid.") from e
    else:
        # If Pillow is not installed, we can't validate, just return the path
        logging.warning("Pillow library not installed. Cannot validate image file, proceeding with path only.")
        return image_path

if __name__ == "__main__":
    import argparse
    import sys

    # Configure basic logging for the CLI tool
    logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')

    parser = argparse.ArgumentParser(description="Extract text from a document.")
    parser.add_argument("file_path", help="Path to the input document file.")

    if len(sys.argv) == 1: # No arguments provided
        parser.print_help(sys.stderr)
        sys.exit(1)

    args = parser.parse_args()

    try:
        text_content, image_refs = parse_input_material(args.file_path)
        if text_content:
            print("\n--- Extracted Text ---")
            print(text_content)
        else:
            print("\nNo text content extracted or file was empty/unsupported.")

        if image_refs:
            print("\n--- Image References ---")
            for ref in image_refs:
                print(ref)
    except FileNotFoundError as e:
        logging.error(e)
        sys.exit(1)
    except Exception as e:
        logging.error(f"An unexpected error occurred: {e}", exc_info=True)
        sys.exit(1) 

